#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py — 全社会議 営業2部月次報告 PowerPoint 生成
出力先: 営業2部/一般薄板課/営業/会議資料/20260421_全社会議_営業2部月次報告.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── カラーパレット ────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x37, 0x5E)   # 濃紺（タイトル背景）
STEEL  = RGBColor(0x2E, 0x6E, 0xA6)   # 鉄鋼ブルー（アクセント）
SILVER = RGBColor(0xD0, 0xD8, 0xE4)   # シルバー（帯）
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE8, 0x6A, 0x1F)   # オレンジ（強調）
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # 薄いブルー背景
GREEN  = RGBColor(0x27, 0xAE, 0x60)

# ─── スライドサイズ: ワイド 16:9 ──────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 完全白紙


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def txb(slide, text, x, y, w, h,
        size=18, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_header_bar(slide, title, subtitle=None):
    """上部ネイビー帯＋タイトル"""
    add_rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    # 左アクセントライン
    add_rect(slide, 0, 0, Inches(0.08), Inches(1.3), fill=ACCENT)

    txb(slide, title,
        Inches(0.2), Inches(0.08), Inches(10), Inches(0.8),
        size=28, bold=True, color=WHITE)

    if subtitle:
        txb(slide, subtitle,
            Inches(0.2), Inches(0.85), Inches(10), Inches(0.4),
            size=13, color=SILVER)

    # 右下 社名
    txb(slide, "三幸商事株式会社 営業2部",
        Inches(9.5), Inches(0.95), Inches(3.6), Inches(0.3),
        size=9, color=SILVER, align=PP_ALIGN.RIGHT)


def add_footer(slide, page_num, total):
    add_rect(slide, 0, Inches(7.15), W, Inches(0.35), fill=NAVY)
    txb(slide, f"{page_num} / {total}",
        Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.32),
        size=9, color=SILVER, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# スライド 1 — タイトル
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# 背景
add_rect(s1, 0, 0, W, H, fill=NAVY)
# 右側グラフィック帯
add_rect(s1, Inches(9.2), 0, Inches(4.13), H, fill=STEEL)
# オレンジ縦ライン
add_rect(s1, Inches(9.1), 0, Inches(0.12), H, fill=ACCENT)

# 会社ロゴ代替テキスト
txb(s1, "三幸商事株式会社",
    Inches(0.5), Inches(0.6), Inches(8), Inches(0.5),
    size=16, color=SILVER)

# メインタイトル
txb(s1, "営業2部 月次報告",
    Inches(0.5), Inches(1.6), Inches(8.5), Inches(1.2),
    size=40, bold=True, color=WHITE)

txb(s1, "2026年4月21日（月）全社会議",
    Inches(0.5), Inches(2.9), Inches(8), Inches(0.6),
    size=20, color=SILVER)

txb(s1, "営業2部 取締役 部長　竹村 翔",
    Inches(0.5), Inches(3.55), Inches(8), Inches(0.55),
    size=18, color=WHITE)

# 右帯テキスト
txb(s1, "守り × 攻め\nの両輪経営",
    Inches(9.35), Inches(2.5), Inches(3.7), Inches(1.8),
    size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txb(s1, "安定収益 × 成長利益",
    Inches(9.35), Inches(4.4), Inches(3.7), Inches(0.5),
    size=14, color=SILVER, align=PP_ALIGN.CENTER)

add_footer(s1, 1, 6)


# ═══════════════════════════════════════════════════════════════
# スライド 2 — 4月度 業績速報
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, W, H, fill=LIGHT)
add_header_bar(s2, "4月度 業績速報", "月中速報（4月中旬時点）")

# 大きな数値カード（左）
def kpi_card(slide, x, y, w, h, label, value, unit, badge, badge_color):
    add_rect(slide, x, y, w, h, fill=WHITE)
    # 左ボーダー
    add_rect(slide, x, y, Inches(0.07), h, fill=badge_color)

    txb(slide, label,
        x + Inches(0.15), y + Inches(0.18), w - Inches(0.2), Inches(0.4),
        size=13, color=STEEL, bold=True)
    txb(slide, value,
        x + Inches(0.15), y + Inches(0.55), w - Inches(0.2), Inches(0.85),
        size=32, bold=True, color=BLACK)
    txb(slide, unit,
        x + Inches(0.15), y + Inches(1.35), w - Inches(0.2), Inches(0.35),
        size=12, color=STEEL)
    # バッジ
    badge_box = slide.shapes.add_shape(1,
        x + w - Inches(0.9), y + Inches(0.15),
        Inches(0.75), Inches(0.45))
    badge_box.fill.solid(); badge_box.fill.fore_color.rgb = badge_color
    badge_box.line.fill.background()
    tb = badge_box.text_frame.paragraphs[0]
    tb.alignment = PP_ALIGN.CENTER
    r = tb.add_run(); r.text = badge
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

kpi_card(s2, Inches(0.4), Inches(1.5), Inches(5.8), Inches(2.1),
         "売  上", "3月超ペース", "月末着地: 1億円超見込み", "◎", GREEN)

kpi_card(s2, Inches(6.5), Inches(1.5), Inches(5.8), Inches(2.1),
         "出荷量", "同上（3月超ペース）", "前年比トレンド維持", "◎", GREEN)

# 月次推移テーブル
table_y = Inches(3.85)
add_rect(s2, Inches(0.4), table_y, Inches(11.9), Inches(0.45), fill=NAVY)
txb(s2, "月次推移",
    Inches(0.5), table_y + Inches(0.05), Inches(4), Inches(0.35),
    size=13, bold=True, color=WHITE)

months = [("1月", "106", "910"), ("2月", "103", "925"), ("3月", "102", "925"), ("4月（速報）", "前月超", "―")]
col_x = [Inches(0.5), Inches(3.3), Inches(6.8), Inches(10.0)]
col_w = [Inches(2.7), Inches(3.3), Inches(3.0), Inches(2.5)]

# ヘッダー行
row_y = table_y + Inches(0.45)
add_rect(s2, Inches(0.4), row_y, Inches(11.9), Inches(0.4), fill=STEEL)
headers = ["月", "売上（百万円）", "出荷量（t）", "判定"]
for i, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    txb(s2, hdr, cx, row_y + Inches(0.05), cw, Inches(0.32),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# データ行
for ri, (mo, uri, deli) in enumerate(months):
    ry = row_y + Inches(0.4) + ri * Inches(0.42)
    bg = LIGHT if ri % 2 == 0 else WHITE
    if ri == 3: bg = RGBColor(0xFF, 0xF3, 0xE0)
    add_rect(s2, Inches(0.4), ry, Inches(11.9), Inches(0.42), fill=bg)

    vals = [mo, uri, deli, "◎" if ri < 3 else "速報"]
    colors = [BLACK, BLACK, BLACK, GREEN if ri < 3 else ACCENT]
    bolds = [ri == 3, ri == 3, ri == 3, True]
    for i, (val, cx, cw, col, bd) in enumerate(zip(vals, col_x, col_w, colors, bolds)):
        txb(s2, val, cx, ry + Inches(0.06), cw, Inches(0.32),
            size=12, bold=bd, color=col, align=PP_ALIGN.CENTER)

txb(s2, "※ 安定推移。現状維持は後退と同じと捉え、2つの成長施策を走らせています。",
    Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
    size=10, color=STEEL, italic=True)

add_footer(s2, 2, 6)


# ═══════════════════════════════════════════════════════════════
# スライド 3 — 成長施策① 缶バッジ事業
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, W, H, fill=LIGHT)
add_header_bar(s3, "【成長施策①】缶バッジ事業の収益化加速", "自社製品 = 価格決定権 = 安定利益")

# 戦略メッセージボックス
add_rect(s3, Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.65), fill=ACCENT)
txb(s3, "鉄鋼（市況依存）× 缶バッジ（自社定価）＝ 収益構造の強化",
    Inches(0.6), Inches(1.52), Inches(12), Inches(0.5),
    size=15, bold=True, color=WHITE)

# 3カラム
col_defs = [
    ("販売体制", STEEL, [
        "新開発マシン（10万〜1,000万円台）の販売本格化",
        "EC開設・在庫5台体制で即納対応",
        "予備機5台によるセンドバック方式でアフター体制を確立",
        "→ 遠方顧客にも対応、販売エリア制約を解消",
    ]),
    ("製造・品質", NAVY, [
        "仕切り導入で不良率を大幅低減",
        "  （1万個ロットで効果実証済み）",
        "全ラインへ横展開予定",
        "板厚 0.20→0.19 のトライ計画中",
        "→ 材料費削減 × 品質向上の両立",
        "金型改定費（約40万円）は投資回収済み",
    ]),
    ("収益インパクト", GREEN, [
        "薄板：安定収益（市況連動）",
        "缶バッジ：成長利益（自社定価）",
        "",
        "この二本柱が会社の",
        "収益構造を強くする",
    ]),
]

col_w3 = Inches(4.0)
col_gap = Inches(0.25)
for ci, (title, color, items) in enumerate(col_defs):
    cx = Inches(0.4) + ci * (col_w3 + col_gap)
    cy = Inches(2.25)
    ch = Inches(4.5)
    add_rect(s3, cx, cy, col_w3, ch, fill=WHITE)
    add_rect(s3, cx, cy, col_w3, Inches(0.5), fill=color)
    txb(s3, title, cx + Inches(0.12), cy + Inches(0.07),
        col_w3 - Inches(0.2), Inches(0.38),
        size=14, bold=True, color=WHITE)

    ty = cy + Inches(0.65)
    for item in items:
        prefix = "• " if item and not item.startswith("→") and not item.startswith(" ") else ""
        icolor = ACCENT if item.startswith("→") else BLACK
        txb(s3, prefix + item, cx + Inches(0.18), ty,
            col_w3 - Inches(0.3), Inches(0.42),
            size=11, color=icolor)
        ty += Inches(0.42)

add_footer(s3, 3, 6)


# ═══════════════════════════════════════════════════════════════
# スライド 4 — 成長施策② TDB新規開拓
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, W, H, fill=LIGHT)
add_header_bar(s4, "【成長施策②】TDB活用による薄板の新規開拓", "依存リスク分散 × 利益率改善")

# 課題ボックス（左）
add_rect(s4, Inches(0.4), Inches(1.5), Inches(5.8), Inches(5.2), fill=WHITE)
add_rect(s4, Inches(0.4), Inches(1.5), Inches(5.8), Inches(0.5), fill=RGBColor(0xC0, 0x39, 0x2B))
txb(s4, "現状の課題", Inches(0.6), Inches(1.57), Inches(5.5), Inches(0.38),
    size=14, bold=True, color=WHITE)

issues = [
    ("TOP10顧客で売上の約50%", "特定先への依存リスク"),
    ("既存深耕だけでは", "売上の天井が見えてきた"),
]
for ri, (iss, sub) in enumerate(issues):
    iy = Inches(2.15) + ri * Inches(1.0)
    add_rect(s4, Inches(0.55), iy, Inches(5.5), Inches(0.75),
             fill=RGBColor(0xFD, 0xED, 0xEB))
    txb(s4, iss, Inches(0.75), iy + Inches(0.05), Inches(5.2), Inches(0.35),
        size=13, bold=True, color=RGBColor(0xC0, 0x39, 0x2B))
    txb(s4, sub, Inches(0.75), iy + Inches(0.38), Inches(5.2), Inches(0.3),
        size=11, color=BLACK)

# 矢印
txb(s4, "→", Inches(6.0), Inches(3.5), Inches(0.8), Inches(0.6),
    size=30, bold=True, color=STEEL, align=PP_ALIGN.CENTER)

# 施策ボックス（右）
add_rect(s4, Inches(6.9), Inches(1.5), Inches(6.0), Inches(5.2), fill=WHITE)
add_rect(s4, Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.5), fill=STEEL)
txb(s4, "TDB活用アプローチ", Inches(7.1), Inches(1.57), Inches(5.7), Inches(0.38),
    size=14, bold=True, color=WHITE)

steps = [
    ("①", "TDB導入", "業種・規模・地域・信用力でスクリーニング"),
    ("②", "DMアプローチ", "ターゲットリストへ一斉送付"),
    ("③", "電話フォロー", "反応先を優先してアポ獲得"),
    ("④", "訪問・商談", "ニーズ確認 → 見積 → 受注"),
]

for si, (num, ttl, desc) in enumerate(steps):
    sy = Inches(2.15) + si * Inches(1.0)
    add_rect(s4, Inches(7.05), sy, Inches(0.5), Inches(0.75), fill=STEEL)
    txb(s4, num, Inches(7.05), sy + Inches(0.18), Inches(0.5), Inches(0.38),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s4, ttl, Inches(7.65), sy + Inches(0.06), Inches(5.0), Inches(0.32),
        size=13, bold=True, color=NAVY)
    txb(s4, desc, Inches(7.65), sy + Inches(0.38), Inches(5.0), Inches(0.3),
        size=11, color=BLACK)

# 期待効果
add_rect(s4, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.55), fill=STEEL)
txb(s4, "期待効果：売上基盤の分散と底上げ ｜ 新規先は単価交渉余地大 → 利益率改善 ｜ 営業2部の攻めの姿勢",
    Inches(0.6), Inches(6.57), Inches(12), Inches(0.4),
    size=12, bold=True, color=WHITE)

add_footer(s4, 4, 6)


# ═══════════════════════════════════════════════════════════════
# スライド 5 — コスト管理
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, W, H, fill=LIGHT)
add_header_bar(s5, "コスト管理", "売上拡大だけでなく、コスト構造の改善にも手を打つ")

cost_items = [
    ("✅ 完了", "物流切替",
     "西濃運輸 → 新潟運輸 直配",
     "年間 約44万円削減",
     GREEN, True),
    ("🔄 検証中", "パレット往復回収",
     "回収ルートの最適化",
     "運賃ゼロ運用を目指す",
     STEEL, False),
    ("📋 策定中", "外材採用検討",
     "内外価格差 約70円/kg",
     "採用基準を策定中",
     ACCENT, False),
]

for ci, (status, title, detail, effect, color, done) in enumerate(cost_items):
    cx = Inches(0.4) + ci * Inches(4.35)
    cy = Inches(1.6)
    cw = Inches(4.1)
    ch = Inches(5.2)

    add_rect(s5, cx, cy, cw, ch, fill=WHITE)
    add_rect(s5, cx, cy, cw, Inches(0.55), fill=color)

    txb(s5, status, cx + Inches(0.12), cy + Inches(0.08),
        cw - Inches(0.2), Inches(0.38),
        size=12, bold=True, color=WHITE)

    txb(s5, title, cx + Inches(0.12), cy + Inches(0.7),
        cw - Inches(0.2), Inches(0.5),
        size=18, bold=True, color=NAVY)

    txb(s5, detail, cx + Inches(0.12), cy + Inches(1.35),
        cw - Inches(0.2), Inches(0.8),
        size=12, color=BLACK)

    # 効果ボックス
    effect_y = cy + Inches(2.4)
    add_rect(s5, cx + Inches(0.12), effect_y, cw - Inches(0.25), Inches(0.75),
             fill=LIGHT)
    txb(s5, "効果", cx + Inches(0.2), effect_y + Inches(0.05),
        Inches(0.6), Inches(0.28), size=9, color=color, bold=True)
    txb(s5, effect, cx + Inches(0.2), effect_y + Inches(0.3),
        cw - Inches(0.4), Inches(0.35),
        size=12, bold=done, color=color)

txb(s5, "売上を伸ばすだけでなく、コスト構造の改善も同時進行で進めています。",
    Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
    size=11, color=NAVY, italic=True, bold=True)

add_footer(s5, 5, 6)


# ═══════════════════════════════════════════════════════════════
# スライド 6 — まとめ
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, W, H, fill=NAVY)
add_rect(s6, 0, 0, Inches(0.12), H, fill=ACCENT)

txb(s6, "まとめ",
    Inches(0.3), Inches(0.2), Inches(12), Inches(0.6),
    size=22, bold=True, color=SILVER)

txb(s6, "営業2部は「守り（安定）× 攻め（成長）」の両輪で動いています。",
    Inches(0.3), Inches(0.85), Inches(12.5), Inches(0.55),
    size=18, bold=True, color=WHITE)

summary_items = [
    ("売  上", "4月は3月超ペース ◎", "月1億円の安定維持", GREEN),
    ("缶バッジ", "販売・品質・アフター体制が整った", "販売件数を積み上げ、収益化を加速", ACCENT),
    ("新規開拓", "TDBでターゲット選定中", "リスト完成 → アプローチ開始", STEEL),
    ("コスト", "物流削減を実行済み（年44万円）", "外材採用基準の策定", RGBColor(0x8E, 0x44, 0xAD)),
]

for ri, (item, current, next_step, color) in enumerate(summary_items):
    ry = Inches(1.6) + ri * Inches(1.25)
    add_rect(s6, Inches(0.3), ry, Inches(12.7), Inches(1.1),
             fill=RGBColor(0x24, 0x4E, 0x7A))

    # カテゴリ
    add_rect(s6, Inches(0.3), ry, Inches(1.8), Inches(1.1), fill=color)
    txb(s6, item, Inches(0.35), ry + Inches(0.3),
        Inches(1.7), Inches(0.5),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 現状
    txb(s6, "現状", Inches(2.2), ry + Inches(0.08),
        Inches(1.0), Inches(0.28), size=9, color=SILVER, bold=True)
    txb(s6, current, Inches(2.2), ry + Inches(0.35),
        Inches(4.5), Inches(0.55),
        size=13, color=WHITE)

    # 区切り
    add_rect(s6, Inches(6.85), ry + Inches(0.15), Inches(0.03), Inches(0.8),
             fill=color)

    # 次の一手
    txb(s6, "次の一手", Inches(7.0), ry + Inches(0.08),
        Inches(1.2), Inches(0.28), size=9, color=SILVER, bold=True)
    txb(s6, next_step, Inches(7.0), ry + Inches(0.35),
        Inches(5.8), Inches(0.55),
        size=13, color=WHITE)

txb(s6, "営業2部 取締役 部長　竹村 翔",
    Inches(9.5), Inches(7.05), Inches(3.7), Inches(0.35),
    size=10, color=SILVER, align=PP_ALIGN.RIGHT)

add_footer(s6, 6, 6)


# ─── 出力 ────────────────────────────────────────────────────
out_dir = os.path.join(os.path.dirname(__file__), "..",
                       "営業2部", "一般薄板課", "営業", "会議資料")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "20260421_全社会議_営業2部月次報告.pptx")
prs.save(out_path)
print(f"✅ 保存完了: {out_path}")
